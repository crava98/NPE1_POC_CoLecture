#!/usr/bin/env python3
"""
Analyze the template to understand its structure.
"""
from pptx import Presentation

template_path = "ppt_templates/modern_corporate_template.potx"

print(f"Analyzing template: {template_path}\n")
print("=" * 60)

prs = Presentation(template_path)

print(f"Available layouts: {len(prs.slide_layouts)}\n")

for i, layout in enumerate(prs.slide_layouts):
    print(f"\n--- Layout [{i}]: {layout.name} ---")
    print(f"Placeholders: {len(layout.placeholders)}")

    for ph in layout.placeholders:
        print(f"  idx={ph.placeholder_format.idx:2d}, type={ph.placeholder_format.type} ({ph.placeholder_format.type})")
        if ph.has_text_frame and ph.text:
            print(f"      Text: {ph.text[:40]}")
