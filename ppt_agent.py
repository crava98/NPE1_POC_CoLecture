"""
Agent 2: PPT Builder Agent
Verantwortlich für die intelligente Erstellung von PowerPoint-Präsentationen.

Der Agent:
1. Ruft verfügbare Templates über MCP ab
2. Analysiert Template-Struktur über MCP
3. Entscheidet intelligente Layout-Wahl basierend auf Content-Typ
4. Generiert professionelle PPT mit optimalen Layouts
"""

import os
import json
import base64
import asyncio
import traceback # Added for detailed exception logging
from io import BytesIO
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_google_genai import ChatGoogleGenerativeAI
from mcp import ClientSession
from mcp.client.sse import sse_client
from data_models import PresentationStructure, ImageColors
from image_providers import get_image_from_gurkli

load_dotenv()
api_key = os.environ.get("GOOGLE_API_KEY")
mcp_server_url = os.environ.get("MCP_SERVER_URL", "http://mcp-server:8000/sse")

# LLM für Layout-Entscheidungen
llm = ChatGoogleGenerativeAI(
    model="gemini-2.5-flash",
    temperature=0.1,  # Niedrig für konsistente Entscheidungen
    google_api_key=api_key
)


import httpx

async def get_templates_from_mcp():
    """Holt die Liste aller verfügbaren Templates vom MCP Server."""
    try:
        async with sse_client(mcp_server_url) as streams:
            async with ClientSession(streams[0], streams[1]) as session:
                await session.initialize()

                print("--> Agent 2: Frage MCP Server nach Templates...")
                result = await session.call_tool("list_templates", arguments={})

                if result.content:
                    templates_json = json.loads(result.content[0].text)
                    return templates_json
    except Exception as e:
        print(f"!!! EXCEPTION in get_templates_from_mcp: {e}")
        traceback.print_exc() # Print the full traceback
        return None
    return None


async def analyze_template_via_mcp(template_name):
    """Analysiert ein Template über den MCP Server."""
    async with sse_client(mcp_server_url) as streams:
        async with ClientSession(streams[0], streams[1]) as session:
            await session.initialize()

            print(f"--> Agent 2: Analysiere Template '{template_name}' via MCP...")
            result = await session.call_tool(
                "analyze_template",
                arguments={"template_name": template_name}
            )

            if result.content:
                analysis = json.loads(result.content[0].text)
                return analysis
    return None


async def get_template_file_from_mcp(template_name):
    """
    Lädt ein Template über MCP (als Base64-encoded Bytes).
    Gibt ein BytesIO-Objekt zurück, das direkt in Presentation() verwendet werden kann.
    """
    async with sse_client(mcp_server_url) as streams:
        async with ClientSession(streams[0], streams[1]) as session:
            await session.initialize()

            print(f"--> Agent 2: Lade Template-Datei '{template_name}' über MCP...")
            result = await session.call_tool(
                "get_template_file",
                arguments={"template_name": template_name}
            )

            if result.content:
                template_data = json.loads(result.content[0].text)

                # Decode Base64 zu Bytes
                template_bytes = base64.b64decode(template_data["data"])

                print(f"  ✓ Template geladen: {template_data['size_mb']} MB")

                # Gib als BytesIO zurück (funktioniert mit Presentation())
                return BytesIO(template_bytes)
    return None


def decide_colors_for_presentation(presentation_data, template_analysis=None):
    """
    Agent entscheidet passende Farben basierend auf dem Präsentationsthema und Template.

    Args:
        presentation_data: PresentationStructure mit allen Slides
        template_analysis: Optional - Analyse des Templates

    Returns:
        dict: {"primary": "#hex", "secondary": "#hex"}
    """
    # Sammle Kontext aus der Präsentation
    all_titles = [slide.title for slide in presentation_data.slides[:5]]  # Erste 5 Titel
    topics_summary = ", ".join(all_titles)

    # Template-Info falls vorhanden
    template_info = ""
    if template_analysis:
        template_info = f"Template: {template_analysis.get('template_name', 'unbekannt')}"

    prompt = f"""
    Du bist ein Experte für Präsentationsdesign und Farbtheorie. Wähle ein passendes Farbschema für diese Präsentation.

    **Präsentationsthemen:**
    {topics_summary}

    {template_info}

    **Richtlinien:**
    - Wähle Farben die zum Thema passen (z.B. Grün für Nachhaltigkeit, Blau für Tech/Business, Orange für Kreativität)
    - Primary: Hauptfarbe für wichtige Elemente
    - Secondary: Komplementär- oder Akzentfarbe
    - Farben sollen professionell und gut lesbar sein
    - Vermeide zu grelle oder schwer kombinierbare Farben

    **Antwort-Format (NUR diese zwei Zeilen, keine Erklärung):**
    primary: #HEXCODE
    secondary: #HEXCODE
    """

    response = llm.invoke(prompt)
    result = response.content.strip()

    # Parse die Antwort
    try:
        lines = result.split("\n")
        primary = "#0066CC"  # Fallback
        secondary = "#00CC66"  # Fallback

        for line in lines:
            line_lower = line.lower().strip()
            if "primary" in line_lower and "#" in line:
                primary = line.split("#")[1][:6]
                primary = f"#{primary}"
            elif "secondary" in line_lower and "#" in line:
                secondary = line.split("#")[1][:6]
                secondary = f"#{secondary}"

        print(f"  🎨 Agent wählt Farben: Primary={primary}, Secondary={secondary}")
        return {"primary": primary, "secondary": secondary}

    except Exception as e:
        print(f"  ⚠ Fehler beim Parsen der Farben: {e} - Verwende Fallback")
        return {"primary": "#0066CC", "secondary": "#00CC66"}


def decide_image_style_for_slide(slide_data):
    """
    Agent entscheidet den besten Bildstil basierend auf dem Slide-Inhalt.

    Returns:
        str: "flat_illustration", "fine_line", oder "photorealistic"
    """
    content_summary = f"Titel: {slide_data.title}\n"
    for item in slide_data.bullets:
        content_summary += f"- {item.bullet}\n"

    prompt = f"""
    Du bist ein Experte für Präsentationsdesign. Wähle den besten Bildstil für diese Folie.

    **Verfügbare Stile:**
    - flat_illustration: Moderne, flache Illustrationen mit klaren Formen und wenig Details.
      IDEAL FÜR: Abstrakte Konzepte, Software/Apps, digitale Transformation, Prozessübersichten, Teamwork-Konzepte, Innovation als Idee.
      BEISPIELE: "Agile Methoden", "Cloud Computing Vorteile", "Unternehmenskultur"

    - fine_line: Feine, detaillierte Linienzeichnungen mit technischem Look.
      IDEAL FÜR: Technische Abläufe, Architektur-Diagramme, wissenschaftliche Konzepte, Blueprints, Engineering.
      BEISPIELE: "Systemarchitektur", "Produktionsprozess", "Technische Spezifikationen"

    - photorealistic: Fotorealistische Bilder mit echten Objekten und Szenen.
      IDEAL FÜR: Physische Produkte, Industrie, Fahrzeuge, Gebäude, Menschen, Natur, reale Wirtschaftsthemen, Handel, Produktion.
      BEISPIELE: "Automobilindustrie", "Handelspolitik", "Produktionsstandorte", "Wirtschaftliche Auswirkungen", "Supply Chain"

    **WICHTIG:**
    - Wenn es um REALE Industrien geht (Auto, Energie, Produktion, Handel) → photorealistic
    - Wenn es um KONZEPTE geht (Strategie, Innovation, digitale Themen) → flat_illustration
    - Wenn es um TECHNISCHE DETAILS geht (Architektur, Prozesse) → fine_line

    **Folie:**
    {content_summary}

    **Deine Aufgabe:**
    Wähle den passendsten Stil. Antworte NUR mit einem der drei Stilnamen: flat_illustration, fine_line, oder photorealistic
    """

    response = llm.invoke(prompt)
    suggested_style = response.content.strip().lower()

    # Validierung
    valid_styles = ["flat_illustration", "fine_line", "photorealistic"]
    if suggested_style not in valid_styles:
        print(f"  ⚠ Ungültiger Stil '{suggested_style}' - Fallback auf flat_illustration")
        return "flat_illustration"

    print(f"  🎨 Agent wählt Bildstil: {suggested_style}")
    return suggested_style


def has_body_placeholder(layout):
    """Prüft ob ein Layout einen echten BODY-Placeholder hat (nicht subTitle)."""
    for ph in layout["placeholders"]:
        ph_type = ph["type"].upper()
        if "BODY" in ph_type and ph["has_text_frame"]:
            return True
    return False


def has_subtitle_only(layout):
    """Prüft ob ein Layout nur subTitle-Placeholders hat (keine BODY)."""
    has_subtitle = False
    for ph in layout["placeholders"]:
        ph_type = ph["type"].upper()
        if "SUB_TITLE" in ph_type or "SUBTITLE" in ph_type:
            has_subtitle = True
        if "BODY" in ph_type:
            return False  # Hat BODY, also nicht "nur subtitle"
    return has_subtitle





def decide_layout_for_slide(template_analysis, slide_data, is_first_slide, slide_index, total_slides):
    """
    Wählt INTELLIGENT das beste Layout für eine Slide.
    Nutzt einen LLM, um basierend auf dem Content den Layout-Typ zu bestimmen.
    """
    layouts = template_analysis["layouts"]
    
    content_summary = f"Titel: {slide_data.title}\n"
    content_summary += f"Punkte: {len(slide_data.bullets)}\n"
    for item in slide_data.bullets:
        content_summary += f"- {item.bullet}\n"
        for sub in item.sub:
            content_summary += f"  - {sub}\n"
    
    # Give the LLM more context about the slide's position
    position_context = ""
    if is_first_slide:
        position_context = "This is the first slide of the presentation."
    elif slide_index == total_slides - 1:
        position_context = "This is the last slide of the presentation."

    layout_categories = [
        "Title and Subtitle",
        "Title and Content",
        "Title, Content and Image",
        "Title Only",
        "Two Content",
        "Image Only",
        "Content Only",
        "Other"
    ]

    prompt = f"""
    You are an expert presentation designer. Your task is to choose the best layout for a slide based on its content.

    **Available Layout Categories:**
    {', '.join(layout_categories)}

    **Content to Classify:**
    - Slide {slide_index + 1} of {total_slides}
    - {position_context}
    - **Title:** {slide_data.title}
    - **Number of Bullets:** {len(slide_data.bullets)}
    - **Content:**
    {content_summary}

    **Examples:**
    - A slide with a title and a few short bullet points should be "Title and Content".
    - A slide with only a title and no bullet points should be "Title Only".
    - The first slide should always be "Title and Subtitle".
    - The last slide, if it's a "Thank you" or "Questions" slide, should be "Title Only".
    - A slide with two main ideas or a comparison should be "Two Content".

    **Your Task:**
    Based on the content and the examples, choose the single most appropriate layout category from the list of available categories.
    Respond ONLY with the name of the category (e.g., "Title and Content").
    """
    
    response = llm.invoke(prompt)
    suggested_category = response.content.strip()
    
    print(f"  LLM-Vorschlag für Content-Typ: {suggested_category}")

    # Matching logic remains the same...
    for layout in layouts:
        if layout.get("classified_type") == suggested_category:
            print(f"  → Layout gefunden (genaue Übereinstimmung mit '{suggested_category}'): {layout['name']}")
            return layout["index"]
            
    # Fallback if no exact match is found
    print(f"  ⚠ Kein exaktes Layout für '{suggested_category}' gefunden. Fallback wird versucht...")
    for layout in layouts:
        if layout.get("classified_type") == "Title and Content":
            print(f"  → Layout gefunden (Fallback auf 'Title and Content'): {layout['name']}")
            return layout["index"]
            
    body_layouts = [l for l in layouts if has_body_placeholder(l)]
    if body_layouts:
        chosen_layout = body_layouts[slide_index % len(body_layouts)]
        print(f"  → Layout gefunden (Fallback auf irgendein Body-Layout): {chosen_layout['name']}")
        return chosen_layout["index"]

    print("  → Absoluter Fallback: Wähle ein einfaches Layout")
    return min(1, len(layouts) - 1)


def generate_ppt_with_agent(presentation_data, language="Deutsch", template_name=None,
                            image_style="flat_illustration", image_mode="auto", image_colors=None):
    """
    Agent 2: Generiert PPT mit intelligenter Layout-Auswahl.

    Args:
        presentation_data: PresentationStructure mit Slides
        language: Sprache der Präsentation
        template_name: Name des Templates (optional)
        image_style: Bildstil (flat_illustration, fine_line, photorealistic)
        image_mode: Bildquelle (auto, stock_only, ai_only)
        image_colors: Farbschema dict {"primary": "#hex", "secondary": "#hex"} (optional)

    Returns:
        Pfad zur generierten PPT
    """
    print("\n" + "="*60)
    print("AGENT 2: PPT BUILDER AGENT")
    print("="*60)

    # 1. Template laden VIA MCP (kein Dateisystem-Zugriff!)
    template_file = None
    template_analysis = None

    if template_name:
        # Hole Template-Datei und Analyse über MCP
        template_file = asyncio.run(get_template_file_from_mcp(template_name))
        template_analysis = asyncio.run(analyze_template_via_mcp(template_name))

        if template_file and template_analysis:
            print(f"✓ Template via MCP geladen: {template_name}")
            print(f"  Dimensionen: {template_analysis['slide_width_inches']}\" x {template_analysis['slide_height_inches']}\"")
            print(f"  Verfügbare Layouts: {template_analysis['total_layouts']}")
            print("  Template Analysis (first 2000 chars):")
            print(json.dumps(template_analysis, indent=2)[:2000])
            if len(json.dumps(template_analysis, indent=2)) > 2000:
                print("  ... (truncated)")
        else:
            print("⚠ Template konnte nicht über MCP geladen werden - verwende Standard")

    # 1b. Farben bestimmen - Agent wählt wenn keine User-Farben
    if image_colors is None:
        print("\n🎨 Keine Farben vorgegeben - Agent wählt passende Farben...")
        image_colors = decide_colors_for_presentation(presentation_data, template_analysis)
    else:
        print(f"\n🎨 User-Farben: Primary={image_colors['primary']}, Secondary={image_colors['secondary']}")

    # 2. PowerPoint erstellen aus Template-Bytes
    if template_file:
        # WICHTIG: Presentation() kann BytesIO direkt verarbeiten!
        prs = Presentation(template_file)
        print(f"  ✓ Presentation aus Template-Bytes erstellt")

        # Lösche Beispiel-Folien aus Template
        slide_count = len(prs.slides)
        for i in range(slide_count - 1, -1, -1):
            rId = prs.slides._sldIdLst[i].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[i]
        print(f"  ✓ {slide_count} Beispiel-Folien entfernt")
    else:
        # Fallback: Standard-Präsentation
        print("  → Erstelle Standard-Präsentation (kein Template)")
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

    # 3. Slides generieren mit intelligenter Layout-Wahl
    print(f"\n{'='*60}")
    print("GENERIERE SLIDES MIT INTELLIGENTER LAYOUT-WAHL")
    print(f"{'='*60}\n")

    total_slides = len(presentation_data.slides)

    for i, slide_data in enumerate(presentation_data.slides):
        print(f"Slide {i+1}: {slide_data.title}")

        # Agent entscheidet Layout
        if template_analysis:
            layout_index = decide_layout_for_slide(
                template_analysis,
                slide_data,
                is_first_slide=(i == 0),
                slide_index=i,
                total_slides=total_slides
            )
            layout = prs.slide_layouts[layout_index]
        else:
            # Ohne Template: Standard-Logik
            layout = prs.slide_layouts[0] if i == 0 else prs.slide_layouts[1]

        slide = prs.slides.add_slide(layout)
        is_title_slide = (i == 0)

        # Titel setzen
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title

        # TITLE SLIDE: Subtitle befüllen (idx=1 ist oft Subtitle)
        if is_title_slide:
            # Erstelle Subtitle aus den ersten Bullets
            subtitle_text = ""
            if slide_data.bullets:
                subtitle_parts = [item.bullet for item in slide_data.bullets[:2]]
                subtitle_text = " | ".join(subtitle_parts)

            # Suche Subtitle-Placeholder (idx=1 bei Title Slides)
            for shape in slide.placeholders:
                if shape.placeholder_format.idx == 1 and shape.has_text_frame:
                    shape.text_frame.text = subtitle_text
                    print(f"  Subtitle gesetzt: '{subtitle_text[:50]}...'")
                    break

        # CONTENT SLIDES: Content einfügen
        else:
            tf = None
            for idx in [1, 2, 3]:  # Versuche verschiedene Indices
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == idx and shape.has_text_frame:
                        tf = shape.text_frame
                        print(f"  Content Placeholder: idx={idx}")
                        break
                if tf:
                    break

            if tf:
                tf.clear()
                for item in slide_data.bullets:
                    p = tf.add_paragraph()
                    p.text = item.bullet
                    p.level = 0
                    for sub in item.sub:
                        ps = tf.add_paragraph()
                        ps.text = sub
                        ps.level = 1
            else:
                print(f"  ⚠ Kein Content-Placeholder gefunden!")

        # Bild einfügen - DYNAMISCH basierend auf Foliengröße!
        if slide_data.unsplashSearchTerms:
            # Setze Bild-Einstellungen auf der Slide (von UI-Parametern)
            # Bei "auto" lässt der Agent den Stil entscheiden
            if image_style == "auto":
                slide_data.style = decide_image_style_for_slide(slide_data)
            else:
                slide_data.style = image_style

            slide_data.image_mode = image_mode
            # Farben sind immer vorhanden (User oder Agent-gewählt)
            slide_data.colors = ImageColors(
                primary=image_colors.get("primary", "#0066CC"),
                secondary=image_colors.get("secondary", "#00CC66")
            )

            image_path = get_image_from_gurkli(slide_data)
            if image_path and os.path.exists(image_path):
                try:
                    # Berechne Position basierend auf Foliengröße
                    slide_width = prs.slide_width
                    slide_height = prs.slide_height

                    # Bild soll ca. 35% der Folienbreite einnehmen
                    img_width = int(slide_width * 0.35)

                    # Position: rechte untere Ecke mit Rand
                    margin = Inches(0.3)
                    img_left = slide_width - img_width - margin
                    img_top = int(slide_height * 0.35)  # Startet bei ca. 35% von oben

                    slide.shapes.add_picture(
                        image_path,
                        left=img_left,
                        top=img_top,
                        width=img_width
                    )
                    print(f"  ✓ Bild eingefügt (Position: {img_left/914400:.1f}\" x {img_top/914400:.1f}\")")
                except Exception as e:
                    print(f"  ⚠ Bild-Fehler: {e}")

        print()

    # 4. Speichern
    output_path = os.path.join("storage", f"generated_presentation_{language}.pptx")
    prs.save(output_path)

    print(f"{'='*60}")
    print(f"✓ AGENT 2: PPT erfolgreich erstellt: {output_path}")
    print(f"{'='*60}\n")

    return output_path
