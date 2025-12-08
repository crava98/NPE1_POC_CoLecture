#!/usr/bin/env python3
"""
Analyze a generated PPTX file to see what went wrong.
"""
import sys
from pptx import Presentation

pptx_path = sys.argv[1] if len(sys.argv) > 1 else "/Users/christophe/Downloads/praesentation_Deutsch (6).pptx"

print(f"Analyzing: {pptx_path}\n")
print("=" * 60)

try:
    prs = Presentation(pptx_path)

    print(f"Slide dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Available layouts: {len(prs.slide_layouts)}\n")

    print("Layout names:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"  [{i}] {layout.name}")

    print("\n" + "=" * 60)
    print("SLIDE DETAILS:\n")

    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        print(f"Layout used: {slide.slide_layout.name}")
        print(f"Number of shapes: {len(slide.shapes)}")

        # Check for title
        if slide.shapes.title:
            print(f"Title: {slide.shapes.title.text[:50]}...")
        else:
            print("Title: [NO TITLE FOUND]")

        # Check for placeholders
        placeholders = [s for s in slide.shapes if s.is_placeholder]
        print(f"Placeholders: {len(placeholders)}")

        for ph in placeholders:
            print(f"  - Placeholder idx={ph.placeholder_format.idx}, type={ph.placeholder_format.type}")

        # Check for images
        images = [s for s in slide.shapes if hasattr(s, 'image')]
        print(f"Images: {len(images)}")

        # Check for text boxes
        textboxes = [s for s in slide.shapes if s.has_text_frame and not s.is_placeholder]
        print(f"Text boxes: {len(textboxes)}")

        # Show all shape types
        print("All shapes:")
        for j, shape in enumerate(slide.shapes):
            shape_type = shape.shape_type
            if shape.has_text_frame:
                text_preview = shape.text[:30].replace('\n', ' ') if shape.text else "[empty]"
                print(f"  [{j}] {shape_type} - Text: {text_preview}")
            else:
                print(f"  [{j}] {shape_type}")

        print()

    print("=" * 60)
    print("ANALYSIS COMPLETE")

except Exception as e:
    print(f"ERROR: {e}")
    import traceback
    traceback.print_exc()
