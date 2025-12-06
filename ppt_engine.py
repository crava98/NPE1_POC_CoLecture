import os
import requests
import json
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt

# Webhook & Token aus .env laden
N8N_WEBHOOK_URL = os.environ.get("N8N_WEBHOOK_URL")
N8N_AUTH_TOKEN = os.environ.get("N8N_AUTH_TOKEN")

def get_image_placeholder(terms):
    """Fallback: Lädt ein Platzhalterbild von loremflickr."""
    search_term = "business"
    if terms and len(terms) > 0:
        search_term = terms[0]
    
    url = f"https://loremflickr.com/800/450/{search_term},work"
    print(f"--> Hole Platzhalter-Bild für: {search_term}")
    
    try:
        response = requests.get(url, timeout=5)
        if response.status_code == 200:
            filename = f"placeholder_{search_term}.jpg"
            path = os.path.join("storage", filename)
            with open(path, "wb") as f:
                f.write(response.content)
            return path
    except Exception:
        pass
    return None

def get_image_from_n8n(slide_data):
    """
    Sendet das KOMPLETTE Slide-Objekt als JSON an n8n.
    Erwartet ein JSON mit einer URL zum Bild als Antwort zurück.
    """
    # 1. Prüfen ob URL da ist
    if not N8N_WEBHOOK_URL:
        print("n8n URL fehlt -> Platzhalter.")
        return get_image_placeholder(slide_data.unsplashSearchTerms)

    # 2. Das Payload bauen (Das volle Objekt!)
    payload = slide_data.model_dump()
    
    # NEU: Korrelations-ID für Tracing hinzufügen
    correlation_id = str(uuid.uuid4())
    payload['correlation_id'] = correlation_id
    print(f"--> Korrelations-ID für n8n-Trace: {correlation_id}")

    # --- DEBUGGING: ZEIGT DIR DAS JSON IM TERMINAL AN ---
    print("\n" + "="*40)
    print(f"--> DEBUG: SENDE DIESES JSON AN N8N:")
    print(json.dumps(payload, indent=2, ensure_ascii=False))
    print("="*40 + "\n")
    # ----------------------------------------------------

    headers = {
        "hslu": N8N_AUTH_TOKEN,
        "Content-Type": "application/json" 
    }

    try:
        # Wir senden JSON und erwarten JSON zurück
        response = requests.post(
            N8N_WEBHOOK_URL,
            data=json.dumps(payload), # Manuell enkodieren
            headers=headers,
            timeout=45
        )

        # 3. Antwort verarbeiten
        if response.status_code == 200:
            print("--> n8n Erfolg! JSON-Antwort empfangen.")
            
            # JSON parsen
            response_data = response.json()
            
            # URL extrahieren
            image_url = response_data[0].get("url")
            if not image_url:
                print("n8n Fehler: Kein 'url' Feld im JSON gefunden.")
                return get_image_placeholder(slide_data.unsplashSearchTerms)

            # Bild von der URL herunterladen
            print(f"--> Lade Bild von URL: {image_url[:50]}...")
            image_response = requests.get(image_url, timeout=30)
            
            if image_response.status_code == 200:
                # Dateinamen basteln
                terms = slide_data.unsplashSearchTerms
                safe_name = "slide_img"
                if terms:
                    safe_name = "".join(x for x in terms[0] if x.isalnum())
                
                filename = f"img_{safe_name}.jpg"
                path = os.path.join("storage", filename)

                # Das Bild speichern (Binary Content)
                with open(path, "wb") as f:
                    f.write(image_response.content)
                return path
            else:
                print(f"Fehler beim Download des Bildes: {image_response.status_code}")
                return get_image_placeholder(slide_data.unsplashSearchTerms)
            
        else:
            print(f"n8n Fehler: {response.status_code} - {response.text}")
            return get_image_placeholder(slide_data.unsplashSearchTerms)

    except Exception as e:
        print(f"n8n Exception: {e}")
        return get_image_placeholder(slide_data.unsplashSearchTerms)

def generate_ppt(presentation_data, language="Deutsch"):
    # Template Logik
    if os.path.exists("template.pptx"):
        prs = Presentation("template.pptx")
    else:
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)

    labels = {
        "Deutsch": "Quellen", "English": "Sources", 
        "Français": "Sources", "Italiano": "Fonti", "Español": "Fuentes"
    }
    source_label = labels.get(language, "Sources")

    for slide_data in presentation_data.slides:
        # Layout Title + Content
        layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        
        # Titel
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.title
        
        # Text Body
        tf = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:
                tf = shape.text_frame
                break
        
        if tf:
            tf.clear()
            if not os.path.exists("template.pptx"):
                tf.width = Inches(8.5)

            for item in slide_data.bullets:
                p = tf.add_paragraph()
                p.text = item.bullet
                p.level = 0
                if p.font: p.font.size = Pt(18)

                for sub in item.sub:
                    ps = tf.add_paragraph()
                    ps.text = sub
                    ps.level = 1
                    if ps.font: ps.font.size = Pt(16)

        # BILD EINFÜGEN (HIER IST DIE ÄNDERUNG)
        if slide_data.unsplashSearchTerms:
            # Wir übergeben jetzt das GANZE slide_data Objekt
            image_path = get_image_from_n8n(slide_data)
            
            if image_path:
                left = Inches(10.0) 
                top = Inches(2.0)
                width = Inches(5.5)
                try:
                    slide.shapes.add_picture(image_path, left, top, width=width)
                except Exception as e:
                    print(f"PPT Bild Fehler: {e}")

        # Quellen
        if slide_data.sources and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            sources_text = f"{source_label}:\n" + "\n".join([f"- {s.documentId} (p. {s.pageNumber})" for s in slide_data.sources])
            text_frame.text = sources_text

    output_path = os.path.join("storage", f"generated_presentation_{language}.pptx")
    prs.save(output_path)
    return output_path