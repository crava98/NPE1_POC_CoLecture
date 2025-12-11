import os
import json
import base64
import re # Added for regex parsing of placeholder types
from pypdf import PdfReader
from pptx import Presentation
from mcp.server import Server
import mcp.types as types
from mcp.server.sse import SseServerTransport
from starlette.applications import Starlette
from starlette.routing import Route
from starlette.requests import Request

# 1. Server definieren
mcp = Server("pdf-and-template-service")
STORAGE_DIR = "/uploads"  # PDF uploads (separate volume)
TEMPLATES_DIR = "/data/templates"  # Templates (baked into image)

# 2. Tool Definition
@mcp.list_tools()
async def list_tools() -> list[types.Tool]:
    return [
        types.Tool(
            name="read_pdf_file",
            description="Liest den Text aus einer PDF-Datei im Speicher.",
            inputSchema={
                "type": "object",
                "properties": {
                    "filename": {"type": "string", "description": "Der Name der Datei (z.B. bericht.pdf)"}
                },
                "required": ["filename"]
            }
        ),
        types.Tool(
            name="list_templates",
            description="Listet alle verfügbaren PowerPoint Templates auf.",
            inputSchema={
                "type": "object",
                "properties": {},
                "required": []
            }
        ),
        types.Tool(
            name="analyze_template",
            description="Analysiert ein PowerPoint Template und gibt Details über verfügbare Layouts, Placeholders und Struktur zurück.",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_name": {"type": "string", "description": "Der Name des Templates (z.B. 'Minimalist Pitch Deck by Slidesgo.pptx')"}
                },
                "required": ["template_name"]
            }
        ),
        types.Tool(
            name="get_template_path",
            description="Gibt den vollständigen Pfad zu einem Template zurück.",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_name": {"type": "string", "description": "Der Name des Templates"}
                },
                "required": ["template_name"]
            }
        ),
        types.Tool(
            name="get_template_file",
            description="Lädt ein Template und gibt es als Base64-encodierten String zurück (für Agent-Zugriff ohne Dateisystem).",
            inputSchema={
                "type": "object",
                "properties": {
                    "template_name": {"type": "string", "description": "Der Name des Templates"}
                },
                "required": ["template_name"]
            }
        )
    ]

# 3. Tool Logik
def classify_layout(layout):
    """
    Analyzes a layout's placeholders to classify it into a generic type.
    """
    print(f"  - Classifying layout: {layout.name}")
    
    # Extract actual placeholder types, ignoring numbers in parentheses
    extracted_placeholder_types = []
    for ph in layout.placeholders:
        ph_type_str = str(ph.placeholder_format.type)
        print(f"    - Raw placeholder type: {ph_type_str}")
        # Regex to extract the type name (e.g., "TITLE", "BODY", "PICTURE")
        match = re.match(r"([A-Z_]+)", ph_type_str)
        if match:
            extracted_placeholder_types.append(match.group(1))
    print(f"    - Extracted placeholder types: {extracted_placeholder_types}")

    has_title = "TITLE" in extracted_placeholder_types or "CENTER_TITLE" in extracted_placeholder_types
    has_body = "BODY" in extracted_placeholder_types
    has_picture = "PICTURE" in extracted_placeholder_types
    has_subtitle = "SUBTITLE" in extracted_placeholder_types

    # Refined classification rules
    if has_title and not has_body and not has_picture:
        if has_subtitle:
            print(f"    - Classified as: Title and Subtitle")
            return "Title and Subtitle"
        print(f"    - Classified as: Title Only")
        return "Title Only"
    elif has_title and has_body and not has_picture:
        print(f"    - Classified as: Title and Content")
        return "Title and Content"
    elif has_title and has_body and has_picture:
        print(f"    - Classified as: Title, Content and Image")
        return "Title, Content and Image"
    elif has_title and has_picture and not has_body:
        print(f"    - Classified as: Title and Image")
        return "Title and Image"
    elif not has_title and has_body and not has_picture:
        print(f"    - Classified as: Content Only")
        return "Content Only"
    elif not has_title and not has_body and has_picture:
        print(f"    - Classified as: Image Only")
        return "Image Only"
    elif extracted_placeholder_types.count("BODY") >= 2:
        print(f"    - Classified as: Two Content")
        return "Two Content"
    elif "VERTICAL_TEXT" in extracted_placeholder_types: # Specific check for this type
        print(f"    - Classified as: Vertical Text")
        return "Vertical Text"
    # Catch-all for layouts with only one placeholder that isn't title/body/picture
    elif len(extracted_placeholder_types) == 1 and not has_title and not has_body and not has_picture:
        print(f"    - Classified as: Single Placeholder")
        return "Single Placeholder" # Could be a chart, table, etc.
    else:
        print(f"    - Classified as: Other (fallback)")
        return "Other" # Default for complex or unhandled layouts

@mcp.call_tool()
async def call_tool(name: str, arguments: dict) -> list[types.TextContent]:
    # PDF Tool
    if name == "read_pdf_file":
        filename = arguments.get("filename")
        file_path = os.path.join(STORAGE_DIR, filename)
        print(f"MCP Server: Lese Datei {file_path}...")

        if not os.path.exists(file_path):
            return [types.TextContent(
                type="text",
                text=f"Fehler: Datei {filename} nicht gefunden. (Pfad geprüft: {file_path})"
            )]

        try:
            reader = PdfReader(file_path)
            text = ""
            for page in reader.pages:
                text += page.extract_text() or ""

            full_text = f"--- INHALT VON {filename} ---\n{text}\n--- ENDE {filename} ---"
            return [types.TextContent(type="text", text=full_text)]

        except Exception as e:
            return [types.TextContent(type="text", text=f"Fehler: {str(e)}")]

    # Template Tools
    elif name == "list_templates":
        print("MCP Server: Liste Templates...")
        try:
            if not os.path.exists(TEMPLATES_DIR):
                return [types.TextContent(type="text", text="Fehler: Templates-Ordner nicht gefunden.")]

            templates = [f for f in os.listdir(TEMPLATES_DIR) if f.endswith(('.pptx', '.potx'))]
            templates.sort()

            result = {
                "count": len(templates),
                "templates": templates
            }
            return [types.TextContent(type="text", text=json.dumps(result, indent=2))]
        except Exception as e:
            return [types.TextContent(type="text", text=f"Fehler: {str(e)}")]

    elif name == "analyze_template":
        template_name = arguments.get("template_name")
        template_path = os.path.join(TEMPLATES_DIR, template_name)
        print(f"MCP Server: Analysiere Template {template_path}...")

        if not os.path.exists(template_path):
            return [types.TextContent(type="text", text=f"Fehler: Template {template_name} nicht gefunden.")]

        try:
            prs = Presentation(template_path)

            # Analysiere Layouts
            layouts_info = []
            for idx, layout in enumerate(prs.slide_layouts):
                placeholders_info = []
                for ph in layout.placeholders:
                    placeholders_info.append({
                        "idx": ph.placeholder_format.idx,
                        "type": str(ph.placeholder_format.type),
                        "has_text_frame": ph.has_text_frame
                    })

                layouts_info.append({
                    "index": idx,
                    "name": layout.name,
                    "classified_type": classify_layout(layout),
                    "placeholders": placeholders_info
                })

            analysis = {
                "template_name": template_name,
                "slide_width_inches": round(prs.slide_width.inches, 2),
                "slide_height_inches": round(prs.slide_height.inches, 2),
                "total_layouts": len(prs.slide_layouts),
                "layouts": layouts_info
            }

            return [types.TextContent(type="text", text=json.dumps(analysis, indent=2))]
        except Exception as e:
            return [types.TextContent(type="text", text=f"Fehler bei Template-Analyse: {str(e)}")]

    elif name == "get_template_path":
        template_name = arguments.get("template_name")
        template_path = os.path.join(TEMPLATES_DIR, template_name)

        if os.path.exists(template_path):
            return [types.TextContent(type="text", text=template_path)]
        else:
            return [types.TextContent(type="text", text=f"Fehler: Template {template_name} nicht gefunden.")]

    elif name == "get_template_file":
        template_name = arguments.get("template_name")
        template_path = os.path.join(TEMPLATES_DIR, template_name)
        print(f"MCP Server: Lade Template-Datei {template_path}...")

        if not os.path.exists(template_path):
            return [types.TextContent(type="text", text=f"Fehler: Template {template_name} nicht gefunden.")]

        try:
            # Lese Template als Bytes und encode als Base64
            with open(template_path, "rb") as f:
                template_bytes = f.read()

            # Base64 encoding für sicheren Transport über JSON/MCP
            template_b64 = base64.b64encode(template_bytes).decode('utf-8')

            # Sende als JSON mit Metadaten
            result = {
                "template_name": template_name,
                "size_bytes": len(template_bytes),
                "size_mb": round(len(template_bytes) / (1024 * 1024), 2),
                "data": template_b64
            }

            print(f"  Template geladen: {result['size_mb']} MB")
            return [types.TextContent(type="text", text=json.dumps(result))]

        except Exception as e:
            return [types.TextContent(type="text", text=f"Fehler beim Laden des Templates: {str(e)}")]

    else:
        raise ValueError(f"Unbekanntes Tool: {name}")

# 4. Web-Server Setup (DER KRITISCHE TEIL)
sse = SseServerTransport("/messages")

async def handle_sse(request: Request):
    # TRICK: Wir geben eine asynchrone Funktion zurück, statt direkt auszuführen.
    # Starlette führt diese Funktion dann mit (scope, receive, send) aus.
    async def asgi_app(scope, receive, send):
        async with sse.connect_sse(scope, receive, send) as streams:
            await mcp.run(streams[0], streams[1], mcp.create_initialization_options())
    
    return asgi_app

async def handle_messages(request: Request):
    # Auch hier: Wir geben die Logik als ASGI-App zurück.
    async def asgi_app(scope, receive, send):
        await sse.handle_post_message(scope, receive, send)
    
    return asgi_app

# 5. App Routes
app = Starlette(routes=[
    Route("/sse", endpoint=handle_sse, methods=["GET"]),
    Route("/messages", endpoint=handle_messages, methods=["POST"])
])