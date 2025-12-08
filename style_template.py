#!/usr/bin/env python3
"""
Enhance an existing PowerPoint template with corporate styling.
This modifies the slide master to add Navy Blue and Accent Blue colors.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Load the existing template
template_path = "ppt_templates/modern_corporate_template.potx"
prs = Presentation(template_path)

# Define colors
COLOR_NAVY_BLUE = RGBColor(0, 51, 102)  # #003366
COLOR_ACCENT_BLUE = RGBColor(230, 240, 255)  # #E6F0FF
COLOR_WHITE = RGBColor(255, 255, 255)

print("Styling template with Modern Corporate colors...")

# Get slide master
slide_master = prs.slide_master

# Set master background to white
try:
    master_bg = slide_master.background
    master_bg.fill.solid()
    master_bg.fill.fore_color.rgb = COLOR_WHITE
    print("✓ Master background set to white")
except Exception as e:
    print(f"⚠ Could not set master background: {e}")

# Add accent bar to slide master
try:
    accent_bar = slide_master.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0),
        top=Inches(0),
        width=Inches(16),
        height=Inches(0.1)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_NAVY_BLUE
    accent_bar.line.fill.background()
    print("✓ Added navy blue accent bar to master")
except Exception as e:
    print(f"⚠ Could not add accent bar: {e}")

# Try to style the Two Content layout specifically
try:
    two_content_layout = next(l for l in prs.slide_layouts if "Two Content" in l.name)

    # Add subtle background accent
    bg_shape = two_content_layout.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=Inches(0),
        top=Inches(0),
        width=Inches(0.15),
        height=Inches(9)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLOR_ACCENT_BLUE
    bg_shape.line.fill.background()
    # Move to back
    two_content_layout.shapes._spTree.remove(bg_shape._element)
    two_content_layout.shapes._spTree.insert(2, bg_shape._element)

    print("✓ Added vertical accent to Two Content layout")
except Exception as e:
    print(f"⚠ Could not style Two Content layout: {e}")

# Save styled template
output_path = "ppt_templates/modern_corporate_styled.potx"
prs.save(output_path)

print(f"\n✓ Styled template saved: {output_path}")
print("\nColors applied:")
print("  - Navy Blue (#003366) for accents")
print("  - Light Blue (#E6F0FF) for subtle backgrounds")
print("  - White (#FFFFFF) for main background")
print("\nThis template will now work with the improved ppt_engine.py!")
print("It will automatically use 'Two Content' layout which supports images.")
